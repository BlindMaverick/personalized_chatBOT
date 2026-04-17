import io
import logging

import pytesseract
from PIL import Image
from PyPDF2 import PageObject, PdfReader

from src.constants import LOG_FILE_PATH
from src.utils import clean_text, setup_logging

# Configure logging
setup_logging()
logger = logging.getLogger(__name__)


def extract_text_from_pdf(file_path: str) -> str:
    """
    Extracts text from a PDF file. Uses OCR if text extraction fails for any page.

    Args:
        file_path (str): Path to the PDF file.

    Returns:
        str: Extracted and cleaned text from the PDF.
    """
    text = ""
    with open(file_path, "rb") as f:
        pdf_reader = PdfReader(f)
        logger.info(f"Opened PDF file for text extraction: {file_path}")

        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            try:
                page_text = page.extract_text()
                if page_text:
                    text += page_text
                    logger.info(f"Extracted text from page {page_num} without OCR.")
                else:
                    logger.info(f"No text found on page {page_num}; attempting OCR.")
                    text += extract_text_from_images(page)
            except Exception as e:
                logger.error(f"Error processing page {page_num}: {e}")

    cleaned_text = clean_text(text)
    logger.info(f"Completed text extraction for {file_path}")
    return cleaned_text


def extract_text_from_images(page: PageObject) -> str:
    """
    Extracts text from images on a page using OCR.

    Args:
        page (PageObject): The PDF page object containing images.

    Returns:
        str: Extracted text from images using OCR.
    """
    text = ""
    for image_file_object in page.images:
        try:
            image = Image.open(io.BytesIO(image_file_object.data))
            ocr_text = pytesseract.image_to_string(image)
            text += ocr_text
            logger.info("Extracted text from image using OCR.")
        except Exception as e:
            logger.error(f"Error processing image for OCR: {e}")
    return text

def extract_text_from_excel(file_path: str) -> str:
    """
    Extracts text from an Excel file.

    Args:
        file_path (str): Path to the Excel file.

    Returns:
        str: Extracted and cleaned text from the Excel file.
    """
    try:
        from openpyxl import load_workbook
    except ModuleNotFoundError as exc:
        raise ModuleNotFoundError(
            "Excel support requires 'openpyxl'. Install it with: pip install openpyxl"
        ) from exc

    text = ""
    workbook = load_workbook(file_path)
    for sheet in workbook.sheetnames:
        ws = workbook[sheet]
        for row in ws.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    text += str(cell) + " "
    cleaned_text = clean_text(text)
    logger.info(f"Extracted text from XLSX file: {file_path}")
    return cleaned_text

def extract_text_from_ppts(file_path: str) -> str:
    """
    Extracts text from a PowerPoint file.

    Args:
        file_path (str): Path to the PowerPoint file.

    Returns:
        str: Extracted and cleaned text from the PowerPoint file.
    """
    try:
        from pptx import Presentation
    except ModuleNotFoundError as exc:
        raise ModuleNotFoundError(
            "PowerPoint support requires 'python-pptx'. Install it with: pip install python-pptx"
        ) from exc

    text = ""
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    cleaned_text = clean_text(text)
    logger.info(f"Extracted text from PPTX file: {file_path}")
    return cleaned_text
